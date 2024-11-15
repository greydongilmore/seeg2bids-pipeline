#!/usr/bin/env python3
# -*- coding: utf-8 -*-

import os,sys
import pandas as pd
import numpy as np
import glob
import math
import datetime
from statistics import NormalDist
import json
import dataframe_image as dfi
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.text import MSO_AUTO_SIZE
from pptx import Presentation
from pptx.util import Inches

if sys.platform=='linux':
	if os.path.exists(os.path.join('/home','greydon','Documents','GitHub','seeg2bids-pipeline')):
		root_dir=os.path.join('/home','greydon','Documents','GitHub','seeg2bids-pipeline')
	else:
		root_dir=os.path.join('/home','greydon','Documents','GitHub','ieegProc')
elif sys.platform=='win32':
	if os.path.exists(os.path.join('C','Users','greydon','Documents','GitHub','seeg2bids-pipeline')):
		root_dir=os.path.join('C','Users','greydon','Documents','GitHub','seeg2bids-pipeline')
	else:
		root_dir=os.path.join('C','Users','greydon','Documents','GitHub','ieegProc')

os.chdir(os.path.join(root_dir,'workflow','scripts','working'))
from helpers import determineFCSVCoordSystem,determine_groups,norm_vec,mag_vec

color_map = {
	"gray": (102, 102, 102),
	"grey": (102, 102, 102),
	"black": (0, 0, 0),
	"red": (255, 0, 0),
	"blue": (42, 96, 153),
	"purple": (128, 0, 128),
	"orange": (255, 128, 0),
	"yellow": (255, 255, 0),
	"brown": (43, 34, 2),
	"green": (0, 169, 51),
	"white": (255, 255, 255),
}

aborted_lang = {
	'skipped',
	'aborted'
}

remap_dict = {
	'Electrode label ("aborted" if skipped)': 'Electrode label',
	'Label (6 chr)': 'Label'
}


def add_slide(presentation, layout, title_dict):
	slide = presentation.slides.add_slide(layout)  # adding a slide
	for ititle in list(title_dict):
		title = slide.shapes.add_textbox(*title_dict[ititle]['position'])
		
		tf = title.text_frame
		tf.auto_size = MSO_AUTO_SIZE.NONE
		tf.word_wrap = False
		
		p = tf.paragraphs[0]
		p.alignment = PP_ALIGN.CENTER
		p.text = ititle
		p.font.size = Pt(title_dict[ititle]['font_size'])
		p.font.color.rgb = title_dict[ititle]['color']

	return slide

def write_cell(table, row, col, value):
	table.cell(row, col).text = "%s" % value

def format_table_header(tbl_tmp):
	
	col_cnt = 0
	header_vals = [
		[(0, col_cnt), (1, 0), 'Electrode']
	]
	
	col_cnt += 1
	if len(tbl_tmp.table.rows[1].cells) == 8:
		header_vals.extend([[(0, col_cnt), (1, 1), 'Implanter']])
		col_cnt += 1
	
	header_vals.extend([
		[(0, col_cnt), (0, col_cnt+1), 'Target Error'],
		[(0, col_cnt+2), (0, col_cnt+3), 'Entry Error'],
		[(0, col_cnt+4), (1, col_cnt+4), 'Radial Angle'],
		[(0, col_cnt+5), (1, col_cnt+5), 'Line Angle'],
		[(1, col_cnt), 'Euclidean'],
		[(1, col_cnt+1), 'Radial'],
		[(1, col_cnt+2), 'Euclidean'],
		[(1, col_cnt+3), 'Radial'],
	])
	
	for ihead in header_vals:
		if len(ihead) > 2:
			cell = tbl_tmp.table.cell(ihead[0][0], ihead[0][1])
			cell.merge(tbl_tmp.table.cell(ihead[1][0], ihead[1][1]))
			cell.text = ihead[2]
		else:
			cell = tbl_tmp.table.cell(ihead[0][0], ihead[0][1])
			cell.text = ihead[1]
			
		cell.vertical_anchor = MSO_ANCHOR.BOTTOM
		cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
		cell.text_frame.paragraphs[0].font.bold = True
		cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
		cell.fill.solid()
		cell.fill.fore_color.rgb = RGBColor(51, 51, 51)
	
	return tbl_tmp

##################################
## Metrics Function Definitions ##
##################################


#Euclidean distance calculation
def euclidianDistanceCalc(xyz_planned, xyz_actual):
	if xyz_planned.ndim>1:
		euc_dist=[]
		for ipoint in range(xyz_planned.shape[0]):
			plan_act_diff = xyz_planned[ipoint] - xyz_actual[ipoint]
			euc_dist.append(math.sqrt(sum(plan_act_diff**2)))
	else:
		plan_act_diff = xyz_planned - xyz_actual
		euc_dist = math.sqrt(sum(plan_act_diff**2))
	return euc_dist


#Radial distance calculation
def radialDistanceCalc(pt, xyz_entry, xyz_target):
	if xyz_entry.ndim>1:
		dist3d=[]
		for ipoint in range(xyz_entry.shape[0]):
			x1_minus_pt = pt[ipoint] - xyz_entry[ipoint]
			x2_minus_x1 = xyz_target[ipoint] - xyz_entry[ipoint]
		
			sumsq_x1_minus_pt = sum(x1_minus_pt * x1_minus_pt)
			sumsq_x2_minus_x1 = sum(x2_minus_x1 * x2_minus_x1)
		
			mydotprod = np.dot(x1_minus_pt, x2_minus_x1)
		
			dist3d.append(np.sqrt((sumsq_x1_minus_pt * sumsq_x2_minus_x1 - (mydotprod * mydotprod))/sumsq_x2_minus_x1))
	else:
		x1_minus_pt = pt - xyz_entry
		x2_minus_x1 = xyz_target - xyz_entry
		
		sumsq_x1_minus_pt = sum(x1_minus_pt * x1_minus_pt)
		sumsq_x2_minus_x1 = sum(x2_minus_x1 * x2_minus_x1)
		
		mydotprod = np.dot(x1_minus_pt, x2_minus_x1)
		
		dist3d = np.sqrt((sumsq_x1_minus_pt * sumsq_x2_minus_x1 - (mydotprod * mydotprod))/sumsq_x2_minus_x1)
	return dist3d


#Radial angle calculation
def ptLineAngleCalc(pt, x_entry, x_target):
	if x_entry.ndim>1:
		deg_angle=[]
		for ipoint in range(x_entry.shape[0]):
			try:
				x1_minus_pt = pt[ipoint] - x_entry[ipoint]
				x2_minus_x1 = x_target[ipoint] - x_entry[ipoint]
			
				sumsq_x1_minus_pt = sum(x1_minus_pt**2)
				sumsq_x2_minus_x1 = sum(x2_minus_x1**2)
			
				mydotprod = np.dot(x1_minus_pt, x2_minus_x1) # sum of products of elements
			
				rad_angle = math.acos(mydotprod/(np.sqrt(sumsq_x1_minus_pt)*np.sqrt(sumsq_x2_minus_x1)))
				deg_angle.append(math.degrees(rad_angle))
			except:
				deg_angle.append(np.nan)
				print(f"Check point {ipoint}")
	else:
		x1_minus_pt = pt - x_entry
		x2_minus_x1 = x_target - x_entry
	
		sumsq_x1_minus_pt = sum(x1_minus_pt**2)
		sumsq_x2_minus_x1 = sum(x2_minus_x1**2)
	
		mydotprod = np.dot(x1_minus_pt, x2_minus_x1) # sum of products of elements
	
		rad_angle = math.acos(mydotprod/(np.sqrt(sumsq_x1_minus_pt)*np.sqrt(sumsq_x2_minus_x1)))
		deg_angle = math.degrees(rad_angle)
	return deg_angle

#Line angle calculation
def lineLineAngleCalc(a_entry, a_target, b_entry, b_target):
	if a_entry.ndim>1:
		deg_angle=[]
		for ipoint in range(a_entry.shape[0]):
			try:
				vectorA = a_target[ipoint] - a_entry[ipoint]
				vectorB = b_target[ipoint] - b_entry[ipoint]
			
				sumsq_vectorA = sum(vectorA**2)
				sumsq_vectorB = sum(vectorB**2)
			
				mydotprod = sum(vectorA*vectorB)
			
				rad_angle = math.acos(mydotprod/(np.sqrt(sumsq_vectorA)*np.sqrt(sumsq_vectorB)))
				deg_angle.append(math.degrees(rad_angle))
			except:
				deg_angle.append(np.nan)
				print(f"Check point {ipoint}")
	else:
		vectorA = a_target - a_entry
		vectorB = b_target - b_entry
	
		sumsq_vectorA = sum(vectorA**2)
		sumsq_vectorB = sum(vectorB**2)
	
		mydotprod = sum(vectorA*vectorB)
	
		rad_angle = math.acos(mydotprod/(np.sqrt(sumsq_vectorA)*np.sqrt(sumsq_vectorB)))
		deg_angle = math.degrees(rad_angle)
	return deg_angle

def mean_confidence_interval(data, confidence=0.95):
	dist = NormalDist.from_samples(data[~np.isnan(data)])
	z = NormalDist().inv_cdf((1 + confidence) / 2.)
	h = dist.stdev * z / ((len(data) - 1) ** .5)
	return dist.mean - h, dist.mean + h



chan_label_dic ={
	3:"RD10R-SP03X",
	4:"RD10R-SP04X",
	5:"RD10R-SP05X",
	6:"RD10R-SP06X",
	7:"RD10R-SP07X"
}


controlpoints_dict={
	"id": "", 	#"vtkMRMLMarkupsFiducialNode_1",
	"label": "",
	"description": "",
	"associatedNodeID": "", #"vtkMRMLScalarVolumeNode1",
	"position": [],
	"orientation": [-1.0, -0.0, -0.0, -0.0, -1.0, -0.0, 0.0, 0.0, 1.0],
	"selected": True,
	"locked": True,
	"visibility": True,
	"positionStatus": "defined"
}


remap_dict={
	'Electrode label ("aborted" if skipped)':'Electrode label',
	'Label (6 chr)':'Label'
}



#%%

debug = True
write_lines = False


if debug:
	class dotdict(dict):
		"""dot.notation access to dictionary attributes"""
		__getattr__ = dict.get
		__setattr__ = dict.__setitem__
		__delattr__ = dict.__delitem__
	
	class Namespace:
		def __init__(self, **kwargs):
			self.__dict__.update(kwargs)
	
	isub='sub-EMOP0471'

	data_dir=r'/home/greydon/Documents/data/emory/derivatives/slicer_scene'
	
	input=dotdict({
				'isub': isub,
				'data_dir':data_dir,
				'shopping_list': os.path.join(data_dir,isub,'*shopping_list.xlsx'),
				'error_metrics': os.path.join(data_dir,isub,f'{isub}_error_metrics.xlsx')
				})
	params=dotdict({
				'sample_line': os.path.join(root_dir,'resources','sample_line.mrk.json'),
				})
	output=dotdict({
		'out_svg': os.path.join(data_dir,isub, f'{isub}_errors.svg'),
		'out_excel': os.path.join(data_dir,isub, f'{isub}_error_metrics.xlsx'),
	})
	
	snakemake = Namespace(output=output, params=params,input=input)

if write_lines:
	with open(snakemake.params.sample_line) as (file):
		sample_line = json.load(file)



isub = snakemake.input.isub
data_dir = snakemake.input.data_dir

patient_files = [x for x in glob.glob(os.path.join(data_dir,isub,'*csv')) if 'space-world' not in os.path.basename(x)]

file_data={}
for ifile in [x for x in patient_files if not x.endswith('empty.csv')]:
	determineFCSVCoordSystem(ifile)
	fcsv_data = pd.read_csv(ifile, skiprows=3, header=None)
	fcsv_data.rename(columns={0:'node_id', 1:'x', 2:'y', 3:'z', 4:'ow', 5:'ox',
			6:'oy', 7:'oz', 8:'vis', 9:'sel', 10:'lock',
			11:'label', 12:'description', 13:'associatedNodeID'}, inplace=True)
	if ifile.lower().endswith('actual.fcsv'):
		file_data['actual']=fcsv_data
	elif ifile.lower().endswith('planned.fcsv'):
		file_data['planned']=fcsv_data
	elif ifile.lower().endswith('seega.fcsv'):
		file_data['seega']=fcsv_data
	elif ifile.lower().endswith('acpc.fcsv'):
		file_data['acpc']=fcsv_data

groupsPlanned, planned_all = determine_groups(np.array(file_data['planned']['label'].values))
label_set=sorted(set(groupsPlanned), key=groupsPlanned.index)

if 'actual' in list(file_data):
	groupsActual, actual_all = determine_groups(np.array(file_data['actual']['label'].values))
	label_set=sorted(set(groupsActual).intersection(groupsPlanned), key=groupsActual.index)
	if not label_set:
		label_set=sorted(set(groupsActual), key=groupsActual.index)

if 'seega' in list(file_data):
	groupsSeega, seega_all = determine_groups(np.array(file_data['seega']['label'].values), True)


shopping_list = glob.glob(f"{os.path.join(data_dir,isub)}/*shopping_list.xlsx")
if shopping_list:
	df_shopping_raw = pd.read_excel(shopping_list[0],header=None)
	df_shopping_list=df_shopping_raw.iloc[4:,:].reset_index(drop=True)
	
	# need to update the column names
	updated_colnames=df_shopping_raw.iloc[3].values
	for idx,ilabel in [(i,x) for i,x in enumerate(updated_colnames) if x in list(remap_dict)]:
		updated_colnames[idx]=remap_dict[ilabel]
	
	df_shopping_list.columns=updated_colnames
	if 'Electrode label' in list(df_shopping_list):
		df_shopping_list=df_shopping_list[df_shopping_list['Electrode label']!='aborted']
	elif 'Serial Num.' in list(df_shopping_list):
		df_shopping_list=df_shopping_list[df_shopping_list['Serial Num.']!='aborted']
	
	df_shopping_list=df_shopping_list.iloc[0:df_shopping_list.loc[:,'Target'].isnull().idxmax()]
	df_shopping_list=df_shopping_list[~df_shopping_list['No.'].isnull()]
	df_shopping_list=df_shopping_list[~df_shopping_list['Target'].isnull()]
	df_shopping_list = df_shopping_list.dropna(axis=1, how='all')
	
	if any(pd.isna(x) for x in list(df_shopping_list)):
		df_shopping_list.drop(np.nan, axis = 1, inplace = True)
	
	if any([x.lower()=='ord.' for x in list(df_shopping_list.keys())]):
		if all(~df_shopping_list.loc[:,'Ord.'].isnull()):
			df_shopping_list=df_shopping_list.sort_values(by=['Ord.']).reset_index(drop=True)
	
	if any([x.lower()=='label' for x in list(df_shopping_list.keys())]):
		df_shopping_list['Label']=[x.strip() for x in df_shopping_list['Label']]
	
	error_idx=[]
	for _,row_elec in df_shopping_list.iterrows():
		if any([x.lower()=='label' for x in list(row_elec.keys())]):
			error_idx.append([i for i,x in enumerate(label_set) if x.lower() == row_elec['Label'].lower().strip()][0])
		else:
			if [i for i,x in enumerate(label_set) if f"({x.lower()})" in row_elec['Target'].lower().strip()]:
				error_idx.append([i for i,x in enumerate(label_set) if f"({x.lower()})" in row_elec['Target'].lower().strip()][0])
			elif [i for i,x in enumerate(label_set) if f"{x.lower()}" in row_elec['Target'].lower().strip()]:
				error_idx.append([i for i,x in enumerate(label_set) if f"{x.lower()}" in row_elec['Target'].lower().strip()][0])
	
	label_set=[label_set[x] for x in error_idx]


mcp_point=None

if 'acpc' in list(file_data):
	ac_point = file_data['acpc'][file_data['acpc']['label'].str.lower() == 'ac'][['x','y','z']].values
	pc_point = file_data['acpc'][file_data['acpc']['label'].str.lower() == 'pc'][['x','y','z']].values
	mcp_point = ((ac_point+pc_point)/2)[0]


elec_data=[]
vtk_cnt=1
for igroup in label_set:
	elec_temp={}
	elec_temp['subject']=isub
	elec_temp['electrode']=igroup
	elec_temp['group']=igroup[1:]
	elec_temp['side']='L' if igroup.startswith('L') else 'R'
	
	if 'seega' in list(file_data):
		seeg_idx=[i for i,x in enumerate(file_data['seega']['label'].values) if igroup in x]
		if all(file_data['seega'].loc[seeg_idx]['description'].isnull()):
			elec_data_temp = file_data['seega'].loc[seeg_idx,['x','y','z']].to_numpy().astype(float)
			dist = np.mean(np.linalg.norm(elec_data_temp[:-1,:] - elec_data_temp[1:,:], axis=1))
			idx_ielec,val_ielec = min(enumerate(list(chan_label_dic)), key=lambda x: abs(x[1]-dist))
			elec_temp['electrodeType']=chan_label_dic[val_ielec]
		else:
			elec_temp['electrodeType']=file_data['seega'].loc[seeg_idx]['description'].mode()[0]
		elec_temp['numContacts']=file_data['seega'].loc[seeg_idx].shape[0]
		
# 		seeg_idx=[i for i,x in enumerate(file_data['seega']['label'].values) if x.startswith(igroup)]
# 		elec_data_temp = file_data['seega'].loc[seeg_idx,['x','y','z']].to_numpy()
# 		dist = np.mean(np.linalg.norm(elec_data_temp[:-1,:] - elec_data_temp[1:,:], axis=1))
# 		idx_ielec,val_ielec = min(enumerate(list(chan_label_dic)), key=lambda x: abs(x[1]-dist))
# 		elec_temp['electrodeType']=chan_label_dic[val_ielec]
# 		elec_temp['numContacts']=file_data['seega'].loc[seeg_idx].shape[0]
	
	if 'planned' in list(file_data):
		planned_idx=[i for i,x in enumerate(file_data['planned']['label'].values) if x.startswith(igroup)]
		elec_temp['plannedTipX']=file_data['planned'].loc[planned_idx,'x'].values[0]
		elec_temp['plannedTipY']=file_data['planned'].loc[planned_idx,'y'].values[0]
		elec_temp['plannedTipZ']=file_data['planned'].loc[planned_idx,'z'].values[0]
		elec_temp['plannedEntryX']=file_data['planned'].loc[planned_idx,'x'].values[1]
		elec_temp['plannedEntryY']=file_data['planned'].loc[planned_idx,'y'].values[1]
		elec_temp['plannedEntryZ']=file_data['planned'].loc[planned_idx,'z'].values[1]
	
	if 'actual' in list(file_data):
		actual_idx=[i for i,x in enumerate(file_data['actual']['label'].values) if x.startswith(igroup)]
		elec_temp['actualTipX']=file_data['actual'].loc[actual_idx,'x'].values[0]
		elec_temp['actualTipY']=file_data['actual'].loc[actual_idx,'y'].values[0]
		elec_temp['actualTipZ']=file_data['actual'].loc[actual_idx,'z'].values[0]
		elec_temp['actualEntryX']=file_data['actual'].loc[actual_idx,'x'].values[1]
		elec_temp['actualEntryY']=file_data['actual'].loc[actual_idx,'y'].values[1]
		elec_temp['actualEntryZ']=file_data['actual'].loc[actual_idx,'z'].values[1]
		
		#need to account for Ad-Tech electrodes encapsulation at the tip. Planned target is electrode tip but
		#actual tip is the edge of the first contact
		mag = mag_vec(file_data['planned'].loc[planned_idx,['x','y','z']].values[0],
		  file_data['planned'].loc[planned_idx,['x','y','z']].values[1])
		norm = norm_vec(file_data['planned'].loc[planned_idx,['x','y','z']].values[0],
		  file_data['planned'].loc[planned_idx,['x','y','z']].values[1])
		plannedTipOffset=file_data['planned'].loc[planned_idx,['x','y','z']].values[1]-(norm*(mag-1))
		
		#elec_temp['plannedOffsetX']=plannedTipOffset[0]
		#elec_temp['plannedOffsetY']=plannedTipOffset[1]
		#elec_temp['plannedOffsetZ']=plannedTipOffset[2]
		
		elec_temp['plannedOffsetX']=elec_temp['plannedTipX']
		elec_temp['plannedOffsetY']=elec_temp['plannedTipY']
		elec_temp['plannedOffsetZ']=elec_temp['plannedTipZ']
		
		xyz_planned_entry = np.array([elec_temp['plannedEntryX'], elec_temp['plannedEntryY'], elec_temp['plannedEntryZ']])
		xyz_actual_entry = np.array([elec_temp['actualEntryX'], elec_temp['actualEntryY'], elec_temp['actualEntryZ']]).T
		xyz_planned_target = np.array([elec_temp['plannedOffsetX'], elec_temp['plannedOffsetY'], elec_temp['plannedOffsetZ']]).T
		xyz_actual_target = np.array([elec_temp['actualTipX'], elec_temp['actualTipY'], elec_temp['actualTipZ']]).T
		
		elec_temp['euclid_dist_target'] = euclidianDistanceCalc(xyz_planned_target, xyz_actual_target)
		elec_temp['euclid_dist_entry'] = euclidianDistanceCalc(xyz_planned_entry, xyz_actual_entry)
		elec_temp['radial_dist_target'] = radialDistanceCalc(xyz_planned_target, xyz_actual_entry, xyz_actual_target)
		elec_temp['radial_dist_entry'] = radialDistanceCalc(xyz_planned_entry, xyz_actual_entry, xyz_actual_target)
		
		if not np.array_equal(np.round(xyz_actual_target,2), np.round(xyz_planned_target,2)):
			try:
				elec_temp['radial_angle'] = ptLineAngleCalc(xyz_actual_target, xyz_planned_entry, xyz_planned_target)
				elec_temp['line_angle'] = lineLineAngleCalc(xyz_actual_entry, xyz_actual_target, xyz_planned_entry, xyz_planned_target)
			except:
				pass
	
	if mcp_point is not None:
		if 'actual' in list(file_data):
			elec_temp['actualTipX_mcp'] = elec_temp['actualTipX']-mcp_point[0]
			elec_temp['actualTipY_mcp']= elec_temp['actualTipY']-mcp_point[1]
			elec_temp['actualTipZ_mcp']= elec_temp['actualTipZ']-mcp_point[2]
			elec_temp['actualEntryX_mcp']= elec_temp['actualEntryX']-mcp_point[0]
			elec_temp['actualEntryY_mcp']= elec_temp['actualEntryY']-mcp_point[1]
			elec_temp['actualEntryZ_mcp']= elec_temp['actualEntryZ']-mcp_point[2]
		if 'planned' in list(file_data):
			elec_temp['plannedTipX_mcp']= elec_temp['plannedOffsetX']-mcp_point[0]
			elec_temp['plannedTipY_mcp']= elec_temp['plannedOffsetY']-mcp_point[1]
			elec_temp['plannedTipZ_mcp']= elec_temp['plannedOffsetZ']-mcp_point[2]
			elec_temp['plannedEntryX_mcp']= elec_temp['plannedEntryX']-mcp_point[0]
			elec_temp['plannedEntryY_mcp']= elec_temp['plannedEntryY']-mcp_point[1]
			elec_temp['plannedEntryZ_mcp']= elec_temp['plannedEntryZ']-mcp_point[2]
	
	if write_lines:
		for itype in ['actual','planned']:
			if itype in list(file_data):
				line_out=sample_line.copy()
				entryp=controlpoints_dict.copy()
				targetp=controlpoints_dict.copy()
				
				entryp["id"]= f"vtkMRMLMarkupsFiducialNode{vtk_cnt}"
				entryp["associatedNodeID"]= f"vtkMRMLScalarVolumeNode{vtk_cnt}"
				entryp["label"]=igroup
				entryp["position"]=[elec_temp[f'{itype}EntryX'],elec_temp[f'{itype}EntryY'],elec_temp[f'{itype}EntryZ']]
				
				targetp["id"]= f"vtkMRMLMarkupsFiducialNode{vtk_cnt+1}"
				targetp["associatedNodeID"]= f"vtkMRMLScalarVolumeNode{vtk_cnt+1}"
				targetp["label"]=igroup
				targetp["position"]=[elec_temp[f'{itype}TipX'],elec_temp[f'{itype}TipY'],elec_temp[f'{itype}TipZ']]
				
				line_out['markups'][0]["controlPoints"]=[targetp,entryp]
				
				json_output = json.dumps(line_out, indent=4)
				with open(os.path.join(data_dir, isub,f'{igroup}_{itype}.mrk.json'), 'w') as (fid):
					fid.write(json_output)
					fid.write('\n')
				
				vtk_cnt+=2
	
	elec_data.append(elec_temp)

elec_data_raw=pd.DataFrame(elec_data)
elec_table=elec_data_raw[['electrode','euclid_dist_target', 'radial_dist_target', 'euclid_dist_entry','radial_dist_entry','radial_angle','line_angle']].round(2)

float_idx=1
# if shopping_list:
# 	if 'implanter' in [x.lower() for x in list(df_shopping_list)]:
# 		float_idx=2
# 		elec_table.insert(1,'implanter', df_shopping_list['Implanter'].values)
# 	

for item in list(elec_table)[float_idx:]:
	elec_table[item]=elec_table[item].astype(float)

if float_idx>1:
	elec_table = elec_table.set_index(['electrode','implanter'])
else:
	elec_table = elec_table.set_index(['electrode'])

elec_table_styled=elec_table.style.map(lambda x: "background-color:#ccffcc;" if x<2 else 'background-color:#ffff00;' if x>=2 and x<3 else "background-color:#ffcccc;")\
	.format('{0:,.2f}').set_properties(**{'text-align': 'center'})

writer = pd.ExcelWriter(snakemake.output.out_excel, engine='openpyxl')
elec_table_styled.to_excel(writer,sheet_name='Sheet1', float_format='%.2f')
writer.close()


pd.set_option('colheader_justify', 'center')
pd.set_option('display.width', -1)
dfi.export(elec_table_styled, snakemake.output.out_svg, table_conversion='firefox')


#%%


pt_pin = 'PIN'
sx_date = 'yyyy-mm-dd'
lastname = "lastname"
firstname = "firstname"

if glob.glob(snakemake.input.shopping_list):
	df_elec_raw = pd.read_excel(
		glob.glob(snakemake.input.shopping_list)[0], header=None)
	df_elec = df_elec_raw.iloc[4:, :].reset_index(drop=True)

	# need to update the column names
	updated_colnames = df_elec_raw.iloc[3].values
	for idx, ilabel in [(i, x) for i, x in enumerate(updated_colnames) if x in list(remap_dict)]:
		updated_colnames[idx] = remap_dict[ilabel]

	df_elec.columns = updated_colnames
	if 'Electrode label' in list(df_elec):
		df_elec = df_elec[df_elec['Electrode label'] != 'aborted']
	elif 'Serial Num.' in list(df_elec):
		df_elec = df_elec[df_elec['Serial Num.'] != 'aborted']

	df_elec = df_elec.iloc[0:df_elec.loc[:, 'Target'].isnull().idxmax()]
	df_elec = df_elec[~df_elec['No.'].isnull()]
	df_elec = df_elec[~df_elec['Target'].isnull()]
	df_elec = df_elec.dropna(axis=1, how='all')

	if any(pd.isna(x) for x in list(df_elec)):
		df_elec.drop(np.nan, axis=1, inplace=True)

	if 'Ord.' in list(df_elec):
		if all(~df_elec.loc[:, 'Ord.'].isnull()):
			df_elec = df_elec.sort_values(by=['Ord.']).reset_index(drop=True)

	pin_idx = [i for i, x in enumerate(
		df_elec_raw.iloc[1].values) if x == 'PIN']
	if pin_idx:
		pt_pin = df_elec_raw.iloc[1, pin_idx[0]+1]

	sx_idx = [i for i, x in enumerate(
		df_elec_raw.iloc[2].values) if x == 'Date']
	if sx_idx:
		if isinstance(df_elec_raw.iloc[2, sx_idx[0]+1], datetime.datetime):
			sx_date = df_elec_raw.iloc[2, sx_idx[0]+1].strftime('%Y-%m-%d')
		elif '_' in df_elec_raw.iloc[2, sx_idx[0]+1]:
			sx_date = datetime.datetime.strptime(
				df_elec_raw.iloc[2, sx_idx[0]+1], '%Y_%m_%d').strftime('%Y-%m-%d')
		else:
			sx_date = datetime.datetime.strptime(
				df_elec_raw.iloc[2, sx_idx[0]+1], '%d/%b/%y').strftime('%Y-%m-%d')

	name_idx = [i for i, x in enumerate(
		df_elec_raw.iloc[0].values) if x == 'Name']
	if name_idx:
		if ',' in df_elec_raw.iloc[0, name_idx[0]+1]:
			lastname, firstname = df_elec_raw.iloc[0, name_idx[0]+1].split(',')
		else:
			firstname, lastname = df_elec_raw.iloc[0, name_idx[0]+1].split(' ')

		firstname = firstname.strip()
		lastname = lastname.strip()

elif os.path.exists(snakemake.input.error_metrics):
	df_elec_raw = pd.read_excel(snakemake.input.error_metrics, header=0)
	df_elec_raw = df_elec_raw.rename(columns={'electrode': 'Electrode label'})
	df_elec = df_elec_raw


prs = Presentation()
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)

blank_slide_layout = prs.slide_layouts[6]
fill = blank_slide_layout.background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(0, 0, 0)

# Title slide
title_dict = {
	f"{lastname}, {firstname}": {
		"font_size": 52,
		"color": RGBColor(255, 255, 255),
		"position": (Inches(3), Inches(2), Inches(10), Inches(1))
	},
	f"{pt_pin}": {
		"font_size": 36,
		"color": RGBColor(255, 255, 255),
		"position": (Inches(5), Inches(3), Inches(6), Inches(.8))
	},
	f"Implantation Date:\n{sx_date}": {
		"font_size": 36,
		"color": RGBColor(255, 255, 255),
		"position": (Inches(5), Inches(4.5), Inches(6), Inches(1.5))
	}
}

title_slide = add_slide(prs, blank_slide_layout, title_dict)
title_slide.name = "title slide"

if glob.glob(snakemake.input.shopping_list):
	# Shopping list
	shopping_list_slide = add_slide(prs, blank_slide_layout, {})
	shopping_list_slide.name = "shopping list"


# Errors
title_dict = {
	"Errors": {
		"font_size": 48,
		"color": RGBColor(255, 255, 255),
		"position": (Inches(3), Inches(.5), Inches(10), Inches(1))
	}
}

errors_data = None
if os.path.exists(snakemake.input.error_metrics):
	errors_data = pd.read_excel(snakemake.input.error_metrics, header=0)
	error_slide = add_slide(prs, prs.slide_layouts[6], title_dict)
	error_slide.name = "errors"
	error_slide.background.fill.solid()
	error_slide.background.fill.fore_color.rgb = RGBColor(0, 0, 0)

	width = Inches(13.0)
	height = Inches(5.0)
	left = (prs.slide_width - width) / 2
	top = (prs.slide_height - height) / 2

	tbl_tmp = error_slide.shapes.add_table(
		errors_data.shape[0]+2, errors_data.shape[1], left, top, width, height)
	tbl = format_table_header(tbl_tmp)

	for row in range(2, len(tbl.table.rows)):
		for cell in range(len(tbl.table.rows[row].cells)):
			if isinstance(errors_data.iloc[row-2, cell], str):
				tbl.table.rows[row].cells[cell].text_frame.text = errors_data.iloc[row-2, cell]
				tbl.table.rows[row].cells[cell].vertical_anchor = MSO_ANCHOR.BOTTOM
				tbl.table.rows[row].cells[cell].text_frame.paragraphs[0].font.bold = True
				tbl.table.rows[row].cells[cell].text_frame.paragraphs[0].font.color.rgb = RGBColor(
					255, 255, 255)
				tbl.table.rows[row].cells[cell].fill.solid()
				tbl.table.rows[row].cells[cell].fill.fore_color.rgb = RGBColor(
					51, 51, 51)
			elif isinstance(errors_data.iloc[row-2, cell], float):
				tbl.table.rows[row].cells[cell].text_frame.text = f"{errors_data.iloc[row-2,cell]:1.2f}"

				if errors_data.iloc[row-2, cell] <= 2:
					col = RGBColor(99, 248, 99)
				elif errors_data.iloc[row-2, cell] > 2 and errors_data.iloc[row-2, cell] < 3:
					col = RGBColor(255, 255, 0)
				else:
					col = RGBColor(255, 95, 54)

				tbl.table.rows[row].cells[cell].fill.solid()
				tbl.table.rows[row].cells[cell].fill.fore_color.rgb = col
				tbl.table.rows[row].cells[cell].fill.fore_color.brightness = 0.4
				tbl.table.rows[row].cells[cell].text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
				tbl.table.rows[row].cells[cell].vertical_anchor = MSO_ANCHOR.BOTTOM

	tbl.left = int((prs.slide_width / 2) - (tbl.width / 2))
	tbl.top = int((prs.slide_height / 2) - (tbl.height / 2))


for _, row_elec in df_elec.iterrows():

	elec_label = None
	if any('Electrode label' in x for x in list(row_elec.keys())):
		elec_label = list(row_elec.keys())[[i for i,x in enumerate(list(row_elec.keys())) if 'Electrode label' in x][0]]
	elif 'Serial Num.' in list(row_elec.keys()):
		elec_label = 'Serial Num.'

	if not any(x in row_elec[list(row_elec.keys())] for x in aborted_lang):

		if any([x.lower() == 'label' for x in list(row_elec.keys())]):
			slide_title = f"{row_elec['Target']} ({row_elec['Label']})"
		elif any([x.lower() == 'target' for x in list(row_elec.keys())]):
			slide_title = row_elec['Target']
		else:
			slide_title = row_elec['Electrode label']

		title_dict = {
			slide_title: {
				"font_size": 48,
				"color": RGBColor(255, 255, 255),
				"position": (Inches(3), Inches(.5), Inches(10), Inches(1))
			}
		}

		elec_slide = add_slide(prs, blank_slide_layout, title_dict)
		elec_slide.name = slide_title

		if errors_data is not None:
			error_idx = []
			if any([x.lower() == 'label' for x in list(row_elec.keys())]):
				error_idx = [i for i, x in enumerate(list(
					errors_data['electrode'].values)) if x.lower() in row_elec['Label'].lower()][0]
			elif any([x.lower() == 'label' for x in list(row_elec.keys())]):
				if [i for i, x in enumerate(errors_data['electrode']) if f'({x.lower()})' in row_elec['Target'].lower()]:
					error_idx = [i for i, x in enumerate(
						errors_data['electrode']) if f'({x.lower()})' in row_elec['Target'].lower()][0]
				elif [i for i, x in enumerate(errors_data['electrode']) if row_elec['Target'].lower().startswith(f'{x.lower()}')]:
					error_idx = [i for i, x in enumerate(
						errors_data['electrode']) if row_elec['Target'].lower().startswith(f'{x.lower()}')][0]
			else:
				if [i for i, x in enumerate(errors_data['electrode']) if f'({x.lower()})' in row_elec['Electrode label'].lower()]:
					error_idx = [i for i, x in enumerate(
						errors_data['electrode']) if f'({x.lower()})' in row_elec['Electrode label'].lower()][0]
				elif [i for i, x in enumerate(errors_data['electrode']) if row_elec['Electrode label'].lower().startswith(f'{x.lower()}')]:
					error_idx = [i for i, x in enumerate(
						errors_data['electrode']) if row_elec['Electrode label'].lower().startswith(f'{x.lower()}')][0]

			if isinstance(error_idx, int):

				width = Inches(13.0)
				height = Inches(1.25)
				left = (prs.slide_width - width) / 2
				top = (prs.slide_height - height) / 10

				tbl = elec_slide.shapes.add_table(
					3, errors_data.shape[1], left, (top*9.5), width, height)
				tbl = format_table_header(tbl)

				cell = 0
				for ival in list(errors_data):
					if isinstance(errors_data.loc[error_idx, ival], str):
						tbl.table.rows[2].cells[cell].text_frame.text = errors_data.loc[error_idx, ival]
						tbl.table.rows[2].cells[cell].vertical_anchor = MSO_ANCHOR.BOTTOM
						tbl.table.rows[2].cells[cell].text_frame.paragraphs[0].font.bold = True
						tbl.table.rows[2].cells[cell].text_frame.paragraphs[0].font.color.rgb = RGBColor(
							255, 255, 255)
						tbl.table.rows[2].cells[cell].fill.solid()
						tbl.table.rows[2].cells[cell].fill.fore_color.rgb = RGBColor(
							51, 51, 51)
					elif isinstance(errors_data.loc[error_idx, ival], float):
						tbl.table.rows[2].cells[cell].text_frame.text = f"{errors_data.loc[error_idx,ival]:1.2f}"

						if errors_data.loc[error_idx, ival] <= 2:
							col = RGBColor(99, 248, 99)
						elif errors_data.loc[error_idx, ival] > 2 and errors_data.loc[error_idx, ival] < 3:
							col = RGBColor(255, 255, 0)
						else:
							col = RGBColor(255, 95, 54)

						tbl.table.rows[2].cells[cell].fill.solid()
						tbl.table.rows[2].cells[cell].fill.fore_color.rgb = col
						tbl.table.rows[2].cells[cell].fill.fore_color.brightness = 0.4
						tbl.table.rows[2].cells[cell].text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
						tbl.table.rows[2].cells[cell].vertical_anchor = MSO_ANCHOR.BOTTOM

					cell += 1

		if isinstance(row_elec[elec_label], int):
			elec_color = 'black'
			elec_text = f"{row_elec[elec_label]}".zfill(3)
		elif '-' in row_elec[elec_label]:
			if row_elec[elec_label].split('-')[0].isdigit():
				elec_color = 'black'
				elec_text = f"{row_elec[elec_label]}".zfill(3)
		elif elec_label == 'Electrode label':
			elec_color = 'black'
			elec_text = f"{row_elec[elec_label]}".zfill(3)
		else:
			elec_color = ''.join(
				[x for x in row_elec[elec_label] if x.isalpha()]).lower()
			elec_text = row_elec[elec_label]

		textbox = elec_slide.shapes.add_textbox(
			Inches(13.5), Inches(4.5), Inches(2), Inches(.5))
		textbox.fill.solid()
		if elec_color in ("yellow", "green", "white"):
			textbox.fill.fore_color.rgb = RGBColor(0, 0, 0)
		else:
			textbox.fill.fore_color.rgb = RGBColor(255, 255, 255)
		tf = textbox.text_frame
		tf.auto_size = MSO_AUTO_SIZE.NONE
		tf.word_wrap = False

		p = tf.paragraphs[0]
		p.alignment = PP_ALIGN.CENTER
		p.text = elec_text
		p.font.size = Pt(24)
		p.font.bold = True
		p.font.color.rgb = RGBColor(
			color_map[elec_color][0], color_map[elec_color][1], color_map[elec_color][2])

		line = textbox.line
		line.color.rgb = RGBColor(255, 0, 0)
		line.width = Inches(0.04)

out_fname = f"{lastname.replace(' ','')}_{firstname}_{sx_date}_maps.pptx"
prs.save(f'{data_dir}/{isub}/{out_fname}')
